#!/usr/bin/env python3
"""
parse.py -- Firecrawl `parse()` equivalent for local files. Free, local, no
API key. Converts PDF/DOCX/XLSX/PPTX/CSV to clean text/markdown so you can
read and reason over it like any other scraped content.

Usage:
  python3 parse.py document.pdf                    # text, page-separated
  python3 parse.py document.pdf --ocr               # OCR fallback for scanned/image PDFs
  python3 parse.py document.docx
  python3 parse.py spreadsheet.xlsx [--sheet "Sheet1"]
  python3 parse.py slides.pptx
  python3 parse.py data.csv

Output: {"path": ..., "type": "pdf"|"docx"|"xlsx"|"pptx"|"csv", "text": "...",
         "pages": N (pdf only), "warning": "..." (if any)}
"""
import argparse
import os
import sys

sys.path.insert(0, __file__.rsplit("/", 1)[0])
from _common import out, eprint


def parse_pdf(path, ocr=False, max_pages=None):
    import fitz  # pymupdf

    doc = fitz.open(path)
    n = len(doc)
    pages_to_read = range(n) if max_pages is None else range(min(max_pages, n))
    texts = []
    used_ocr_pages = []

    for i in pages_to_read:
        page = doc[i]
        text = page.get_text().strip()
        if not text and ocr:
            try:
                import pytesseract
                from PIL import Image
                import io

                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = pytesseract.image_to_string(img).strip()
                if text:
                    used_ocr_pages.append(i + 1)
            except ImportError:
                text = ""
                eprint("[parse] OCR requested but pytesseract/PIL not installed "
                       "(pip install --break-system-packages pytesseract pillow, "
                       "plus the tesseract binary itself, e.g. `brew install tesseract`)")
        texts.append(f"--- page {i + 1} ---\n{text}")

    result = {"path": path, "type": "pdf", "pages": n, "text": "\n\n".join(texts)}
    empty_pages = sum(1 for t in texts if t.split("\n", 1)[-1].strip() == "")
    if empty_pages:
        result["warning"] = (
            f"{empty_pages}/{len(texts)} page(s) had no extractable text — likely "
            f"scanned/image pages. Re-run with --ocr" +
            (" (tesseract not installed, see stderr)" if not used_ocr_pages and ocr else ".")
        )
    if used_ocr_pages:
        result["ocr_pages"] = used_ocr_pages
    return result


def parse_docx(path):
    from docx import Document

    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            style = (p.style.name or "").lower()
            prefix = "## " if "heading" in style else ""
            parts.append(prefix + p.text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if rows:
            parts.append("\n".join(rows))
    return {"path": path, "type": "docx", "text": "\n\n".join(parts)}


def parse_xlsx(path, sheet=None):
    import pandas as pd

    xl = pd.ExcelFile(path)
    sheets = [sheet] if sheet else xl.sheet_names
    parts = []
    for name in sheets:
        df = xl.parse(name)
        parts.append(f"--- sheet: {name} ---\n{df.to_markdown(index=False)}")
    return {"path": path, "type": "xlsx", "sheets": xl.sheet_names, "text": "\n\n".join(parts)}


def parse_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"--- slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append(" | ".join(c.text.strip() for c in row.cells))
        parts.append("\n".join(lines))
    return {"path": path, "type": "pptx", "slides": len(prs.slides), "text": "\n\n".join(parts)}


def parse_csv(path):
    import pandas as pd

    df = pd.read_csv(path)
    return {"path": path, "type": "csv", "rows": len(df), "columns": list(df.columns),
            "text": df.to_markdown(index=False)}


EXT_MAP = {
    ".pdf": parse_pdf, ".docx": parse_docx, ".xlsx": parse_xlsx,
    ".xls": parse_xlsx, ".pptx": parse_pptx, ".csv": parse_csv,
}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--ocr", action="store_true", help="PDF only: OCR pages with no extractable text")
    ap.add_argument("--sheet", default=None, help="XLSX only: parse a single sheet by name")
    ap.add_argument("--max-pages", type=int, default=None, help="PDF only: cap pages read")
    args = ap.parse_args()

    if not os.path.exists(args.path):
        out({"error": f"file not found: {args.path}"})
        sys.exit(1)

    ext = os.path.splitext(args.path)[1].lower()
    handler = EXT_MAP.get(ext)
    if not handler:
        out({"error": f"unsupported file type '{ext}'. Supported: {sorted(EXT_MAP)}"})
        sys.exit(1)

    try:
        if handler is parse_pdf:
            result = parse_pdf(args.path, ocr=args.ocr, max_pages=args.max_pages)
        elif handler is parse_xlsx:
            result = parse_xlsx(args.path, sheet=args.sheet)
        else:
            result = handler(args.path)
    except ImportError as e:
        out({"error": f"missing dependency for '{ext}' files: {e}. "
                       f"pip install --break-system-packages pymupdf python-docx pandas openpyxl python-pptx tabulate"})
        sys.exit(1)

    out(result)


if __name__ == "__main__":
    main()
